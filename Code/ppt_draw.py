from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from lxml import etree
import json
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from math import cos, sin, pi
from Set_Text import TextRunFactory
from utils import textwrap, requests
from config import ollama_url  # 移除 used_model



# 因為pptx本身沒有L形狀的連線方式，所以自創
def add_l_connector(slide, start, bend, end): 
    """ 插入 L 型連線（兩段）並組成群組，尾端加上箭頭 """
    left1, top1 = start[0], start[1]
    left2, top2 = bend[0], bend[1]
    left3, top3 = end[0], end[1]

    # 插入兩段 connector：start → bend → end
    line1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left1, top1, left2, top2)
    line2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left3, top3, left2, top2)

    for line in [line1, line2]:
        line.line.width = Pt(2)
        line.line.color.rgb = RGBColor(0, 0, 0)
        line.line.end_arrowhead = None  # 先清掉預設箭頭

    # 只在尾端 line2 加箭頭（XML 操作）
    line_elem = line2.line._get_or_add_ln()
    line_elem.append(parse_xml(
        '<a:headEnd type="triangle" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
    ))

    # 將兩段線組成一個群組
    group = slide.shapes.add_group_shape([line1, line2])
    return group

# 把每個節點的前後順序連接起來(判斷最不會擋住節點的連線方式)
def draw_connectors(slide, nodes, shapes, layout_key):
    for node_id, node in nodes.items():
        start_shape = shapes[node_id]
        for nxt_id in node["next"]:
            if nxt_id not in shapes:
                continue
            end_shape = shapes[nxt_id]

            # 基本幾何定位點
            start_top_y = start_shape.top
            start_center_x = start_shape.left + start_shape.width / 2
            start_center_y = start_shape.top + start_shape.height / 2
            start_bottom_y = start_shape.top + start_shape.height
            start_right_x = start_shape.left + start_shape.width
            start_left_x = start_shape.left
            
            end_top_y = end_shape.top
            end_center_x = end_shape.left + end_shape.width / 2
            end_center_y = end_shape.top + end_shape.height / 2
            end_left_x = end_shape.left
            end_right_x = end_shape.left + end_shape.width

            dx = end_shape.left - start_shape.left
            dy = end_shape.top - start_shape.top

            arrow_margin = int(start_shape.width/5)

            # 條件 1：A 在 B 正上方（x 座標差小，y 座標差為正）
            if abs(dx) < 20 and dy > 0:
                # ↓ 從 A 底部 → B 頂部
                start_x = int(start_center_x)
                start_y = int(start_bottom_y)
                end_x = int(end_center_x)
                end_y = int(end_top_y)
                
            # 條件 2：A 在 B 正左方（y 差小，x 差為正）
            elif abs(dy) < 20 and dx > 0:
                # → 從 A 右側 → B 左側
                start_x = int(start_right_x)
                start_y = int(start_center_y)
                end_x = int(end_left_x)
                end_y = int(end_center_y)
			
			# 條件 4：A 在 B 右側（新增）
            elif abs(dy) < 20 and dx < 0:
                # ← 從 A 左側 → B 右側
                start_x = int(start_left_x)
                start_y = int(start_center_y)
                end_x = int(end_right_x)
                end_y = int(end_center_y)

            # 條件 3：A 在 B 左上方，且 B 左邊 > A 中心（代表 B 稍往右偏）
            elif dx > 0 and dy > 0 and end_left_x > start_center_x:
                # ↘ 從 A 底部 → B 左側
                start_x = int(start_center_x)+arrow_margin
                start_y = int(start_bottom_y)
                end_x = int(end_left_x)
                end_y = int(end_center_y)
                print(node_id, nxt_id, "右下")
                add_l_connector(slide, [start_x, start_y], [start_x, end_y], [end_x, end_y])
                continue
            # 條件 6：A 在 B 左下方
            elif dx > 0 and dy < 0 :
                # 從 A 上方 → B 左側
                start_x = int(start_center_x)+arrow_margin
                start_y = int(start_top_y)
                end_x = int(end_left_x)
                end_y = int(end_center_y)
                print(node_id, nxt_id, "右上")
                add_l_connector(slide, [start_x, start_y], [start_x, end_y], [end_x, end_y])
                continue
			# 條件 5：A 在 B 右上方
            elif dx < 0 and dy > 0 :
                # ↙ 從 A 底部 → B 右側
                start_x = int(start_center_x)-arrow_margin
                start_y = int(start_bottom_y)
                end_x = int(end_right_x)
                end_y = int(end_center_y)
                add_l_connector(slide, [start_x, start_y], [start_x, end_y], [end_x, end_y])
                continue
            # 條件 7：A 在 B 右下方
            elif dx < 0 and dy < 0 :
                # ↙ 從 A 底部 → B 右側
                start_x = int(start_center_x)-arrow_margin
                start_y = int(start_top_y)
                end_x = int(end_right_x)
                end_y = int(end_center_y)
                add_l_connector(slide, [start_x, start_y], [start_x, end_y], [end_x, end_y])
                continue
            else:
                # 預設 fallback：從 A 右側 → B 左側
                start_x = int(start_right_x)
                start_y = int(start_center_y)
                end_x = int(end_left_x)
                end_y = int(end_center_y)
                print(node_id, nxt_id, "一般")

            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.ELBOW, end_x, end_y, start_x, start_y
            )
            connector.line.color.rgb = RGBColor(0, 0, 0)
            connector.line.width = Pt(2)
            line_elem = connector.line._get_or_add_ln()
            line_elem.append(parse_xml(
                '<a:headEnd type="triangle" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            ))

# 檢查兩個框框是否重疊，保留 margin 間距
def is_overlap(rect1, rect2, margin=10):
    x1, y1, w1, h1 = rect1
    x2, y2, w2, h2 = rect2

    return not (
        x1 + w1 + margin <= x2 or  # rect1 在 rect2 左邊 + margin
        x2 + w2 + margin <= x1 or  # rect2 在 rect1 左邊 + margin
        y1 + h1 + margin <= y2 or  # rect1 在 rect2 上方 + margin
        y2 + h2 + margin <= y1     # rect2 在 rect1 上方 + margin
    )

# 建立投影片與避免重疊的自動排版
def create_slide(prs, nodes, layout_key, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    if not title_shape:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.7))
    title_shape.text = title_text

    shapes = {}
    shape_positions = []  # 儲存所有已擺放的框框位置 (x, y, w, h)
    H_START = 100

    for node_id, node in nodes.items():
        layout = node["layouts"][layout_key]
        x = Inches(layout["x"] * 1.4 / 100)
        y = Inches((H_START + layout["y"] * 1.6) / 100)
        w = Inches(layout["width"] * 1.2 / 100)
        h = Inches(layout["height"] * 1.6 / 100)

        # 嘗試避免與現有框框重疊
        max_attempts = 100
        dy = Inches(0.15)  # 每次微調的距離
        attempt = 0
        while any(is_overlap((x, y, w, h), rect) for rect in shape_positions) and attempt < max_attempts:
            y += dy
            attempt += 1

        # 加入框框
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(253, 234, 218)
        shape.line.color.rgb = RGBColor(0, 0, 0)

        # 文字樣式
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.clear()

        run_factory = TextRunFactory(p)
        run_factory.add_icon(node.get("icon", ""), font_size=35, color=RGBColor(152, 72, 7))
        run_factory.add_id(node["id"])

        shapes[node_id] = shape
        shape_positions.append((x, y, w, h))  # 記錄此框框位置

    draw_connectors(slide, nodes, shapes, layout_key)

def create_list_slide(prs, nodes, title_text="清單圖"):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    if not title_shape:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.7))
    title_shape.text = title_text

    shapes = {}

    num_nodes = len(nodes)
    spacing = Inches(0.2)
    total_width = prs.slide_width - Inches(0.6)
    shape_width = (total_width - spacing * (num_nodes - 1)) / num_nodes
    shape_height = Inches(1)
    top_margin = Inches(2)

    for i, (node_id, node) in enumerate(nodes.items()):
        left = Inches(0.3) + i * (shape_width + spacing)

        shape_1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_margin, shape_width, shape_height)
        shape_1.fill.background()  # 讓填色透明
        shape_1.line.fill.background()  # 讓邊框也透明（完全不顯示）
        tf_1 = shape_1.text_frame
        tf_1.clear()

        # 第一行 icon
        p_1 = tf_1.paragraphs[0]
        p_1.alignment = PP_ALIGN.CENTER
        p_1.clear()

        run_factory = TextRunFactory(p_1)
        run_factory.add_icon(node.get("icon", ""))

        shapes[node_id] = shape_1

        text_len = len(node["id"]+node["add"])
        if text_len < 50:
            shape_2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_margin + shape_height, shape_width, shape_height*2)
        else:
            shape_2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top_margin + shape_height, shape_width, shape_height*3)
        # 設定填色與邊框
        shape_2.fill.solid()
        shape_2.fill.fore_color.rgb = RGBColor(218, 235, 247)
        shape_2.line.color.rgb = RGBColor(0, 0, 0)

        tf_2 = shape_2.text_frame
        tf_2.clear()

        p_2 = tf_2.paragraphs[0]
        p_2.alignment = PP_ALIGN.CENTER
        p_2.space_after = Pt(4)
        p_2.clear()

        # 第二行 ID, 第三行 Add
        run_factory_2 = TextRunFactory(p_2)
        run_factory_2.add_id(node["id"])
        run_factory_2.add_add(node["add"])

def create_cycle_slide(prs, nodes, title_text="循環圖"):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    if not title_shape:
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.7))
    title_shape.text = title_text

    shapes = {}

    # 圓的參數
    num_nodes = len(nodes)
    radius = Inches(1.5 + num_nodes * 0.1)
    center_x = prs.slide_width / 2
    center_y = prs.slide_height / 2 + Inches(0.5)

    shape_w = Inches(2)
    shape_h = Inches(1.8)

    for i, (node_id, node) in enumerate(nodes.items()):
        angle = 2 * pi * i / num_nodes - pi / 2  # 每個節點的角度
        x = center_x + radius * cos(angle) - shape_w / 2
        y = center_y + radius * sin(angle) - shape_h / 2

        # 圓上的框框
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, shape_w, shape_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(221, 235, 247)
        shape.line.color.rgb = RGBColor(0, 0, 0)

        tf = shape.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.clear()

        # 第一行 icon, 第二行 ID, 第三行 Add
        run_factory = TextRunFactory(p)
        run_factory.add_icon(node.get("icon", ""))
        run_factory.add_id(node["id"])

        

def detect_layout_types(nodes):
    for node in nodes.values():
        return list(node.get("layouts", {}).keys())[:2]
    return []

# 修正：加入 model 與 temperature 參數
def create_node(summary_text, model="gpt-oss:20b", temperature=0.7):
    prompt = textwrap.dedent(f"""
    請根據以下的格式，萃取出其中的主要流程或架構，並將其拆解為結構化的節點信息。每個節點應該包含以下欄位：
    - **id**: 節點的唯一識別字串，這將是該節點的顯示文字（不重複，適當時可以使用中文）。
    - **add**: 此節點的簡要說明, 最多不超過20字。
    - **next**: 這是一個列表，包含當前節點的指向節點（即後續步驟）。如果此節點是流程的最終步驟，則設為空列表 []。
    - **icon**: 這是與節點相關的 Unicode 表情符號，用來表示該節點的象徵意圖或功能。

    請注意：
    - 若無明確結構或層級，請合理推測分層並找出關聯性。
    - 根據流程描述，拆解成 2 到 4 個主要節點，並在這些節點之間建立「下一步」的關聯。
    - 保證每個節點都包含以上所列的欄位。

    ### 範例：
    {{
        "節點1": {{
            "id": "自動化核心概念",
            "add": "",
            "next": ["觸發節點", "動作節點"],
            "icon": "🔍"
        }},
        "節點2": {{
            "id": "觸發節點",
            "add": "何時執行工作流程？",
            "next": ["資料儲存"],
            "icon": "🧩"
        }}
    }}

    ### 以下是要處理的文字：
    {summary_text}
    """)
    
    # 修改：透過 payload 傳遞動態的 model 與 temperature
    payload = {
        "model": model,
        "prompt": prompt,
        "stream": False,
        "options": {
            "temperature": temperature
        }
    }
    
    response = requests.post(f"{ollama_url}/api/generate", json=payload)

    if response.status_code == 200:
        text = response.json().get("response", "")
        if text:
            # 1. 移除可能的 HTML/XML 標籤 (例如 DeepSeek 常見的 <think> 標籤)
            text = re.sub(r"<.*?>", "", text, flags=re.DOTALL)
            
            # 2. 用正則表達式尋找最外層的 {}
            m = re.search(r'(\{[\s\S]*\})', text)
            json_str = m.group(1) if m else text
            
            try:
                nodes = json.loads(json_str)
                return nodes
            except json.JSONDecodeError:
                print("Error: JSON 解讀失敗，原始回傳內容為：", text)
                return None
    return None

def generate_diagram_to_ppt(save_path, st_status, node_data):
    if not node_data:
        msg = "⚠️ 節點資料為空或解析失敗，略過圖表生成。"
        print(msg)
        if st_status:
            st_status.warning(msg)
        return
    
    print(node_data)
    msg = "📊 製作圖表中..."
    if st_status:
        st_status.info(msg)

    prs = Presentation(save_path)

    nodes = {n["id"]: n for n in node_data.values()}  # 正確轉換 key

    layout_types = detect_layout_types(nodes)
    for layout in layout_types:
        create_slide(prs, nodes, layout, layout)
    create_list_slide(prs, nodes, "清單圖")
    create_cycle_slide(prs, nodes, "循環圖")

    msg = "📊 製作完成!"
    if st_status:
        st_status.info(msg)
    prs.save(save_path)      # 儲存為 PPT 檔案