from utils import Presentation, json, requests, os, textwrap, re, RGBColor, PP_ALIGN
from config import ollama_url
from pptx.util import Pt

# 先對每一個code進行分析用途, 以及結構
def call_llm_individual_file(file_content, model, temperature):
    prompt = textwrap.dedent(f"""
    You are a professional Python engineer and technical writer. 
    Based on the following content, please analyze and describe its purpose and structure clearly.

    The content may include code, plain text, documentation, or mixed formats. 
    If it is code, explain its functionality and structure.
    If it is text or documentation, summarize its main points and logical flow.

    Do not include any extra formatting or code block markers.
    Only output plain text.

    The content is as follows:
    {file_content}
    """)
    
    payload = {
        "model": model,
        "prompt": prompt,
        "stream": False,
        "options": {
            "temperature": temperature  # Ollama 的參數通常放在 options 內
        }
    }
    response = requests.post(f"{ollama_url}/api/generate", json=payload)
    return response.json().get("response", "").strip()

# 再針對第一個LLM所給的code分析去做summary, 並轉成"title", "content"的json格式, 以便轉成powerpoint輸出
def call_llm_summary(all_summaries_text, num_pages=5, level="專家", language="中文", model="gpt-oss:20b", temperature=0.7):
    prompt = textwrap.dedent(f"""
    You are a professional presentation assistant. Based on the following summary, please create a presentation with exactly {num_pages} slides.

    - Target audience level: {level}
    - Output language: {language}
    - Each slide must contain a "title" and "multi-line content"
    - The result must strictly follow the JSON format below, without any additional explanations or extra text
    - Only output the content body, no slide numbers or section headings

    Example format (please follow this structure exactly):
    {{
    "slides": [
        {{
        "title": "System Overview",
        "content": "The system consists of multiple modules\\nSupports vector search, file upload, and reranking"
        }},
        {{
        "title": "Process Steps",
        "content": "1. User uploads a file\\n2. Convert to vector\\n3. Search and rerank"
        }}
    ]
    }}

    Now, based on the summary below, please generate the JSON format output:
    ===== Summary Start =====
    {all_summaries_text}
    ===== Summary End =====
    """)
    
    payload = {
        "model": model,
        "prompt": prompt,
        "stream": False,
        "options": {
            "temperature": temperature
        }
    }

    response = requests.post(f"{ollama_url}/api/generate", json=payload)
    text = response.json().get("response", "").strip()

    # print("===== 第二層回應原始內容 =====")
    # print(text)

    text = re.sub(r"<.*?>", "", text)   # 移除可能殘留的 HTML 標籤（如 <p>、<think>）

    # 嘗試擷取 JSON 區塊
    m = re.search(r'(\{[\s\S]*\})', text)
    json_str = m.group(1) if m else text

    try:
        return json.loads(json_str)     # 嘗試解析為 JSON 格式
    except json.JSONDecodeError:
        print("JSON 解讀失敗，內容為：", json_str)
        return {"slides": []}

# 負責呼叫第一個和第二個LLM, 並彙整為簡報結構
def generate_report(all_texts, st_status, num_pages, level, language, model="gpt-oss:20b", temperature=0.7):
    
    summaries = []
    # 第一階段：分析每個檔案
    for i, content in enumerate(all_texts):
        msg = f"▶ 分析第 {i+1} 個檔案中 (使用模型: {model})..."
        if st_status:
            st_status.info(msg)
        
        # 傳遞模型與溫度
        summary = call_llm_individual_file(content, model, temperature)
        summaries.append(f"【第{i+1}個檔案】\n{summary}")

    summary_text = "\n\n---\n\n".join(summaries)

    # 第二階段：彙整簡報
    msg = f"📊 正在彙整簡報結構 (使用模型: {model})..."
    if st_status:
        st_status.info(msg)
        
    # 傳遞模型與溫度
    structure = call_llm_summary(
        summary_text, 
        num_pages=num_pages, 
        level=level, 
        language=language, 
        model=model, 
        temperature=temperature
    )

    if st_status:
        st_status.info("📊 彙整完成!")
    return structure

# 　把簡報 JSON 結構轉為真正的 .pptx 簡報
def generate_ppt_from_report(structure_json, save_path, font_size=18):
    os.makedirs(os.path.dirname(save_path), exist_ok=True)

    prs = Presentation()

    # 如果預設簡報已有空白頁，先移除
    if prs.slides:
        r_id = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(r_id)

    layout = prs.slide_layouts[1]   # 使用標準標題+內容的版型

    # 為每一頁簡報建立 Slide
    for slide_data in structure_json.get("slides", []):
        title = slide_data.get("title", "無標題")
        content = slide_data.get("content", "")

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.text = ""  # 清空預設文字

        # 每行內容逐行加入簡報段落
        for idx, line in enumerate(content.split("\n")):
            if idx == 0:
                tf.text = line
                p = tf.paragraphs[0]  # 取得第一個段落
            else:
                p = tf.add_paragraph()
                p.text = line
            p.font.size = Pt(font_size)     # 調整字體大小
    prs.save(save_path)      # 儲存為 PPT 檔案
